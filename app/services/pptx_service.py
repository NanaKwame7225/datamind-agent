"""
DataMind Agent — PowerPoint Export Service
Generates professional .pptx slide decks from analysis results.
One slide per section: title, executive summary, charts, findings, recommendations, finance.
"""
from __future__ import annotations
import io, base64, logging, re
from datetime import datetime, timezone

logger = logging.getLogger(__name__)

# Brand palette
TEAL = (0x0a, 0x4d, 0x4a)
ACCENT = (0x00, 0xc8, 0xbe)
DARK = (0x0b, 0x12, 0x21)
LIGHT = (0xdc, 0xe8, 0xf5)
GREY = (0x66, 0x66, 0x66)
RED = (0xf0, 0x40, 0x60)
WARN = (0xf0, 0xa0, 0x20)
GREEN = (0x00, 0xd8, 0x88)


class PowerPointExportService:
    """Builds a polished PowerPoint deck from a DataMind analysis result."""

    def build_deck(self, result: dict, finance: dict = None, chart_images: list = None) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        industry = result.get("industry", "General").replace("_", " ").title()
        query = result.get("query", "Data Analysis")
        date_str = datetime.now(timezone.utc).strftime("%d %B %Y")

        # ── SLIDE 1: TITLE ──
        s = prs.slides.add_slide(blank)
        self._fill_bg(s, prs, DARK)
        self._text(s, "DataMind Agent", Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.1),
                   44, ACCENT, bold=True)
        self._text(s, f"{industry} Analysis Report", Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.8),
                   26, LIGHT, bold=True)
        self._text(s, f'"{query}"', Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.9),
                   15, (0x8a, 0xa0, 0xc0), italic=True)
        self._text(s, f"{date_str}  •  NkaySolutions  •  Accra, Ghana", Inches(0.9), Inches(6.5),
                   Inches(11.5), Inches(0.5), 12, GREY)

        # ── SLIDE 2: EXECUTIVE SUMMARY ──
        narrative = result.get("narrative", "")
        exec_summary = self._extract_section(narrative, "EXECUTIVE SUMMARY")
        if exec_summary:
            s = prs.slides.add_slide(blank)
            self._slide_header(s, prs, "Executive Summary")
            self._text(s, exec_summary, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.2),
                       16, (0x22, 0x33, 0x44), line_spacing=1.4)

        # ── SLIDE 3+: KEY METRICS ──
        metrics = result.get("metrics", [])
        if metrics:
            s = prs.slides.add_slide(blank)
            self._slide_header(s, prs, "Key Metrics")
            n = min(len(metrics), 4)
            card_w = Inches(2.7)
            gap = Inches(0.35)
            total_w = card_w * n + gap * (n - 1)
            start_x = (prs.slide_width - total_w) / 2
            for i, m in enumerate(metrics[:4]):
                x = start_x + i * (card_w + gap)
                self._metric_card(s, x, Inches(2.6), card_w, Inches(2.2), m)

        # ── CHART SLIDES ──
        if chart_images:
            for ci, img_data in enumerate(chart_images[:6]):
                if not img_data or not img_data.get("image"):
                    continue
                s = prs.slides.add_slide(blank)
                self._slide_header(s, prs, img_data.get("title", f"Chart {ci+1}"))
                try:
                    img_bytes = base64.b64decode(img_data["image"].split(",")[-1])
                    img_stream = io.BytesIO(img_bytes)
                    pic_w = Inches(10)
                    pic_x = (prs.slide_width - pic_w) / 2
                    s.shapes.add_picture(img_stream, pic_x, Inches(1.7), width=pic_w)
                except Exception as e:
                    logger.warning(f"Chart image {ci} failed: {e}")
                if img_data.get("subtitle"):
                    self._text(s, img_data["subtitle"], Inches(0.9), Inches(6.7), Inches(11.5),
                               Inches(0.5), 12, GREY, italic=True, align=PP_ALIGN.CENTER)

        # ── FINDINGS SLIDE ──
        insights = result.get("insights", [])
        if insights:
            s = prs.slides.add_slide(blank)
            self._slide_header(s, prs, "Key Findings")
            y = Inches(1.7)
            for i in insights[:5]:
                sev = i.get("severity", "info")
                color = {"critical": RED, "warning": WARN, "success": GREEN}.get(sev, ACCENT)
                conf = f"  ({int(i.get('confidence',0)*100)}% confidence)" if i.get('confidence') else ""
                self._bullet_row(s, y, i.get("title", ""), i.get("body", ""), color, conf)
                y += Inches(1.0)

        # ── RECOMMENDATIONS SLIDE ──
        recs = self._extract_recommendations(narrative)
        if recs:
            s = prs.slides.add_slide(blank)
            self._slide_header(s, prs, "Recommendations")
            y = Inches(1.7)
            for idx, rec in enumerate(recs[:6], 1):
                self._numbered_row(s, y, idx, rec)
                y += Inches(0.85)

        # ── FINANCE SLIDES ──
        if finance:
            self._add_finance_slides(prs, blank, finance)

        # ── CLOSING SLIDE ──
        s = prs.slides.add_slide(blank)
        self._fill_bg(s, prs, DARK)
        self._text(s, "Thank You", Inches(0.9), Inches(2.8), Inches(11.5), Inches(1),
                   40, ACCENT, bold=True, align=PP_ALIGN.CENTER)
        self._text(s, "Generated by DataMind Agent — AI Data Analysis Platform",
                   Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.6), 14, LIGHT, align=PP_ALIGN.CENTER)
        self._text(s, "NkaySolutions  •  Accra, Ghana", Inches(0.9), Inches(4.6), Inches(11.5),
                   Inches(0.5), 12, GREY, align=PP_ALIGN.CENTER)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    # ── HELPERS ──
    def _fill_bg(self, slide, prs, rgb):
        from pptx.dml.color import RGBColor
        from pptx.util import Inches
        shp = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(*rgb)
        shp.line.fill.background()
        shp.shadow.inherit = False
        slide.shapes._spTree.remove(shp._element)
        slide.shapes._spTree.insert(2, shp._element)

    def _slide_header(self, slide, prs, title):
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        # Accent bar
        bar = slide.shapes.add_shape(1, Inches(0.9), Inches(0.7), Inches(0.15), Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*ACCENT); bar.line.fill.background()
        self._text(slide, title, Inches(1.2), Inches(0.6), Inches(11), Inches(0.8), 28, TEAL, bold=True)

    def _text(self, slide, text, x, y, w, h, size, rgb, bold=False, italic=False,
              align=None, line_spacing=None):
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align: p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = RGBColor(*rgb)
        run.font.name = 'Calibri'
        return tb

    def _metric_card(self, slide, x, y, w, h, metric):
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        card = slide.shapes.add_shape(1, x, y, w, h)
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xf0, 0xf4, 0xf8)
        card.line.color.rgb = RGBColor(*ACCENT); card.line.width = Pt(1)
        card.shadow.inherit = False
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_top = Inches(0.2); tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(metric.get("label","")).upper()
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = RGBColor(*GREY)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = str(metric.get("value",""))
        r2.font.size = Pt(30); r2.font.bold = True; r2.font.color.rgb = RGBColor(*TEAL)
        if metric.get("change_pct") is not None:
            p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
            cp = metric.get("change_pct")
            r3 = p3.add_run(); r3.text = f"{'+' if cp>0 else ''}{cp}% vs previous"
            r3.font.size = Pt(11)
            r3.font.color.rgb = RGBColor(*(GREEN if cp>0 else RED if cp<0 else GREY))

    def _bullet_row(self, slide, y, title, body, color, conf=""):
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        dot = slide.shapes.add_shape(1, Inches(0.9), y+Inches(0.08), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(*color); dot.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(1.3), y, Inches(11.3), Inches(0.95))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title + conf
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = RGBColor(0x1a, 0x20, 0x30)
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = body[:180]
        r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(*GREY)

    def _numbered_row(self, slide, y, num, text):
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        circ = slide.shapes.add_shape(9, Inches(0.9), y, Inches(0.5), Inches(0.5))
        circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(*ACCENT); circ.line.fill.background()
        ctf = circ.text_frame; ctf.word_wrap = False
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = str(num); cr.font.size = Pt(16); cr.font.bold = True
        cr.font.color.rgb = RGBColor(0, 0, 0)
        tb = slide.shapes.add_textbox(Inches(1.6), y, Inches(11), Inches(0.75))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)[:200]
        r.font.size = Pt(14); r.font.color.rgb = RGBColor(0x1a, 0x20, 0x30)

    def _add_finance_slides(self, prs, blank, finance):
        from pptx.util import Inches
        tax = finance.get("tax"); acct = finance.get("accounting"); fraud = finance.get("fraud")
        if tax and not tax.get("error"):
            s = prs.slides.add_slide(blank)
            self._slide_header(s, prs, "Tax Analysis")
            metrics = tax.get("metrics", [])
            n = min(len(metrics), 4)
            if n:
                card_w = Inches(2.7); gap = Inches(0.35)
                total_w = card_w*n + gap*(n-1); start_x = (prs.slide_width-total_w)/2
                for i,m in enumerate(metrics[:4]):
                    self._metric_card(s, start_x+i*(card_w+gap), Inches(2.0), card_w, Inches(2.0), m)
            y = Inches(4.4)
            for f in tax.get("findings", [])[:3]:
                sev=f.get("severity","info"); color={"critical":RED,"warning":WARN}.get(sev,ACCENT)
                self._bullet_row(s, y, f.get("title",""), f.get("body",""), color); y+=Inches(0.9)
        if acct and not acct.get("error"):
            s = prs.slides.add_slide(blank)
            self._slide_header(s, prs, "Accounting Analysis")
            score = acct.get("health_score", 0)
            grade = 'A' if score>=85 else 'B' if score>=70 else 'C' if score>=55 else 'D' if score>=40 else 'F'
            self._text(s, f"Balance Sheet Health: {grade}  ({score}/100)", Inches(0.9), Inches(1.7),
                       Inches(11.5), Inches(0.7), 20, TEAL, bold=True)
            metrics = acct.get("metrics", [])
            n = min(len(metrics), 4)
            if n:
                card_w=Inches(2.7); gap=Inches(0.35); total_w=card_w*n+gap*(n-1); start_x=(prs.slide_width-total_w)/2
                for i,m in enumerate(metrics[:4]):
                    self._metric_card(s, start_x+i*(card_w+gap), Inches(2.7), card_w, Inches(2.0), m)
        if fraud and not fraud.get("error"):
            s = prs.slides.add_slide(blank)
            self._slide_header(s, prs, "Fraud Detection")
            score = fraud.get("risk_score", 0); level = fraud.get("risk_level", "Low")
            color = RED if score>=60 else WARN if score>=20 else GREEN
            self._text(s, f"Risk Level: {level}  ({score}/100)", Inches(0.9), Inches(1.8),
                       Inches(11.5), Inches(0.8), 24, color, bold=True)
            y = Inches(3.0)
            for f in fraud.get("findings", [])[:4]:
                sev=f.get("severity","info"); c={"critical":RED,"warning":WARN}.get(sev,ACCENT)
                self._bullet_row(s, y, f.get("title",""), f.get("body",""), c); y+=Inches(0.9)
        actions = finance.get("priority_actions", [])
        if actions:
            s = prs.slides.add_slide(blank)
            self._slide_header(s, prs, "Priority Actions")
            y = Inches(1.7)
            for idx,a in enumerate(actions[:6],1):
                self._numbered_row(s, y, idx, f"[{a.get('module','')}] {a.get('action','')}")
                y += Inches(0.85)

    def _extract_section(self, narrative, section_name):
        if not narrative: return ""
        narrative = re.sub(r'^\s*\{"narrative"\s*:\s*"', '', narrative)
        narrative = narrative.replace('\\n', '\n').replace('\\"', '"')
        pattern = rf'#{{1,3}}\s*{re.escape(section_name)}\s*\n(.*?)(?=\n#{{1,3}}\s|\Z)'
        m = re.search(pattern, narrative, re.DOTALL | re.IGNORECASE)
        if m:
            text = m.group(1).strip()
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            text = re.sub(r'\n+', ' ', text)
            return text[:600]
        return ""

    def _extract_recommendations(self, narrative):
        if not narrative: return []
        narrative = narrative.replace('\\n', '\n')
        recs = []
        m = re.search(r'#{1,3}\s*RECOMMENDATIONS(.*?)(?=\n#{1,3}\s|\Z)', narrative, re.DOTALL | re.IGNORECASE)
        block = m.group(1) if m else narrative
        for line in block.split('\n'):
            line = line.strip()
            if re.match(r'^\d+\.', line):
                clean = re.sub(r'^\d+\.\s*', '', line)
                clean = re.sub(r'\*\*(.+?)\*\*', r'\1', clean)
                if len(clean) > 8:
                    recs.append(clean)
        return recs[:6]


pptx_service = PowerPointExportService()
